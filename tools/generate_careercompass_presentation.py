#!/usr/bin/env python3
"""Generate the formal CareerCompass graduation-project presentation."""

from __future__ import annotations

import argparse
import math
import tempfile
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORT = ROOT / "Career_Compass_Submission.docx"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Aptos"
FONT_DISPLAY = "Aptos Display"

NAVY = "15324B"
NAVY_2 = "234A68"
TEAL = "0E6A5F"
TEAL_2 = "16877B"
TEAL_PALE = "E6F2F0"
BLUE = "3478B8"
BLUE_PALE = "EAF2F8"
INK = "101820"
MUTED = "5C6A78"
WHITE = "FFFFFF"
CANVAS = "F4F6F8"
CARD = "FFFFFF"
BORDER = "DDE3E9"
SHADOW = "E6EBF0"
GREEN = "18794E"
GREEN_PALE = "E8F5EE"
AMBER = "8A5300"
AMBER_PALE = "FFF4DE"
RED = "B42318"
RED_PALE = "FEECEB"
SLATE_PALE = "EDF1F5"


def color(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def set_slide_background(slide, fill_color: str = CANVAS) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color(fill_color)


def style_run(run, size: float, font_color: str = INK, bold: bool = False, name: str = FONT) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color(font_color)
    run.font.bold = bold


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    font_color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.0,
    font_name: str = FONT,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    if line_spacing is not None:
        paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    style_run(run, size, font_color, bold, font_name)
    return box


def add_rich_line(slide, parts, x, y, w, h, *, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0)
    frame.margin_top = frame.margin_bottom = Inches(0)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for part in parts:
        run = paragraph.add_run()
        run.text = part[0]
        style_run(
            run,
            part[1] if len(part) > 1 else 16,
            part[2] if len(part) > 2 else INK,
            part[3] if len(part) > 3 else False,
            part[4] if len(part) > 4 else FONT,
        )
    return box


def add_bullets(
    slide,
    items: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    font_color: str = INK,
    bullet_color: str = TEAL,
    gap: float = 7,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0)
    frame.margin_top = frame.margin_bottom = Inches(0)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.05
        bullet_run = paragraph.add_run()
        bullet_run.text = "•  "
        style_run(bullet_run, size, bullet_color, True)
        text_run = paragraph.add_run()
        text_run.text = item
        style_run(text_run, size, font_color, False)
    return box


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = BORDER,
    line_width: float = 1.0,
):
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(line)
    shape.line.width = Pt(line_width)
    return shape


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = CARD,
    line: str = BORDER,
    shadow: bool = True,
):
    if shadow:
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + 0.045,
            y + 0.055,
            w,
            h,
            fill=SHADOW,
            line=SHADOW,
            line_width=0.3,
        )
    return add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=fill,
        line=line,
        line_width=0.85,
    )


def add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float = 0.34,
    *,
    fill: str = TEAL_PALE,
    font_color: str = TEAL,
    size: float = 10.5,
    bold: bool = True,
):
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=fill,
        line=fill,
        line_width=0.2,
    )
    return add_text(
        slide,
        text,
        x,
        y + 0.005,
        w,
        h - 0.01,
        size=size,
        font_color=font_color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_number_badge(slide, number: str, x: float, y: float, size: float = 0.52, fill: str = TEAL):
    add_shape(slide, MSO_SHAPE.OVAL, x, y, size, size, fill=fill, line=fill, line_width=0.3)
    add_text(
        slide,
        number,
        x,
        y,
        size,
        size,
        size=14,
        font_color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide, slide_number: int, total_main: int = 14, section: str = "CAREERCOMPASS") -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.62),
        Inches(7.13),
        Inches(12.71),
        Inches(7.13),
    )
    line.line.color.rgb = color(BORDER)
    line.line.width = Pt(0.8)
    add_text(slide, section, 0.64, 7.18, 2.2, 0.18, size=8.5, font_color=MUTED, bold=True)
    add_text(
        slide,
        f"{slide_number:02d}",
        12.13,
        7.16,
        0.55,
        0.2,
        size=9,
        font_color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    progress = min(slide_number / total_main, 1.0)
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        0.0,
        7.44,
        SLIDE_W,
        0.06,
        fill=BORDER,
        line=BORDER,
        line_width=0.1,
    )
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        0.0,
        7.44,
        SLIDE_W * progress,
        0.06,
        fill=TEAL,
        line=TEAL,
        line_width=0.1,
    )


def add_slide_title(slide, number: int, title: str, eyebrow: str | None = None, *, appendix: bool = False) -> None:
    set_slide_background(slide)
    if eyebrow:
        add_text(
            slide,
            eyebrow.upper(),
            0.66,
            0.30,
            5.6,
            0.24,
            size=10.5,
            font_color=TEAL,
            bold=True,
        )
    add_text(
        slide,
        title,
        0.64,
        0.56,
        11.9,
        0.58,
        size=27,
        font_color=NAVY,
        bold=True,
        font_name=FONT_DISPLAY,
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.64, 1.18, 0.72, 0.055, fill=TEAL, line=TEAL, line_width=0.1)
    add_footer(slide, number, section="TECHNICAL APPENDIX" if appendix else "CAREERCOMPASS")


def add_image_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = w / h
    if ratio >= box_ratio:
        image_w = w
        image_h = w / ratio
        image_x = x
        image_y = y + (h - image_h) / 2
    else:
        image_h = h
        image_w = h * ratio
        image_x = x + (w - image_w) / 2
        image_y = y
    return slide.shapes.add_picture(
        str(path), Inches(image_x), Inches(image_y), Inches(image_w), Inches(image_h)
    )


def add_image_cover(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    box_ratio = w / h
    picture = slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if image_ratio > box_ratio:
        crop = (1 - box_ratio / image_ratio) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < box_ratio:
        crop = (1 - image_ratio / box_ratio) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def add_picture_frame(slide, x: float, y: float, w: float, h: float, line: str = BORDER) -> None:
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame.fill.background()
    frame.line.color.rgb = color(line)
    frame.line.width = Pt(0.9)


def add_right_arrow(slide, x: float, y: float, w: float, h: float, fill: str = TEAL) -> None:
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, x, y, w, h, fill=fill, line=fill, line_width=0.2)


def add_compass(slide, x: float, y: float, size: float, *, muted: bool = False) -> None:
    stroke = "80A7A1" if muted else TEAL
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    outer.fill.background()
    outer.line.color.rgb = color(stroke)
    outer.line.width = Pt(2.2)
    inner = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + size * 0.18),
        Inches(y + size * 0.18),
        Inches(size * 0.64),
        Inches(size * 0.64),
    )
    inner.fill.background()
    inner.line.color.rgb = color(BORDER)
    inner.line.width = Pt(1.0)
    vertical = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x + size / 2),
        Inches(y + size * 0.13),
        Inches(x + size / 2),
        Inches(y + size * 0.87),
    )
    vertical.line.color.rgb = color(stroke)
    vertical.line.width = Pt(1.0)
    horizontal = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x + size * 0.13),
        Inches(y + size / 2),
        Inches(x + size * 0.87),
        Inches(y + size / 2),
    )
    horizontal.line.color.rgb = color(stroke)
    horizontal.line.width = Pt(1.0)
    north = add_shape(
        slide,
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        x + size * 0.39,
        y + size * 0.13,
        size * 0.22,
        size * 0.45,
        fill=TEAL,
        line=TEAL,
        line_width=0.2,
    )
    north.rotation = 0
    south = add_shape(
        slide,
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        x + size * 0.39,
        y + size * 0.42,
        size * 0.22,
        size * 0.45,
        fill=NAVY,
        line=NAVY,
        line_width=0.2,
    )
    south.rotation = 180
    add_shape(
        slide,
        MSO_SHAPE.OVAL,
        x + size * 0.45,
        y + size * 0.45,
        size * 0.10,
        size * 0.10,
        fill=WHITE,
        line=NAVY,
        line_width=0.8,
    )


def add_metric_card(slide, value: str, label: str, x: float, y: float, w: float, *, note: str = "") -> None:
    add_card(slide, x, y, w, 1.12, fill=WHITE, shadow=False)
    add_text(slide, value, x + 0.18, y + 0.13, w - 0.36, 0.38, size=23, font_color=NAVY, bold=True)
    add_text(slide, label, x + 0.18, y + 0.53, w - 0.36, 0.25, size=11.5, font_color=TEAL, bold=True)
    if note:
        add_text(slide, note, x + 0.18, y + 0.79, w - 0.36, 0.20, size=9.5, font_color=MUTED)


def add_two_line_card(slide, title: str, body: str, x: float, y: float, w: float, h: float, *, accent: str = TEAL, fill: str = WHITE) -> None:
    add_card(slide, x, y, w, h, fill=fill)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.07, h, fill=accent, line=accent, line_width=0.1)
    add_text(slide, title, x + 0.94, y + 0.16, w - 1.14, 0.30, size=15.5, font_color=NAVY, bold=True)
    add_text(slide, body, x + 0.94, y + 0.53, w - 1.14, h - 0.66, size=12.5, font_color=MUTED, line_spacing=1.02)


def extract_assets(report: Path, output_dir: Path) -> dict[str, Path]:
    media = {
        "meu": "word/media/image27.jpg",
        "transcript": "word/media/image34.png",
        "dashboard": "word/media/image21.png",
        "mentors": "word/media/image1.png",
    }
    extracted: dict[str, Path] = {}
    with zipfile.ZipFile(report, "r") as archive:
        for key, member in media.items():
            destination = output_dir / Path(member).name
            destination.write_bytes(archive.read(member))
            extracted[key] = destination
    extracted["gap_algorithm"] = ROOT / "docs/figures/chapter6/algorithm-3-skill-gap-priority.png"
    extracted["skill_vector"] = ROOT / "docs/figures/chapter6/algorithm-2-student-skill-vector.png"
    extracted["course_ranking"] = ROOT / "docs/figures/chapter6/algorithm-4-course-recommendation-ranking.png"
    extracted["quiz"] = ROOT / "docs/figures/chapter6/algorithm-5-quiz-generation-grading.png"
    extracted["architecture"] = ROOT / "docs/figures/architecture.png"
    return extracted


def make_deck(output_path: Path) -> None:
    if not REPORT.exists():
        raise FileNotFoundError(f"Required report not found: {REPORT}")

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    blank_layout = presentation.slide_layouts[6]
    presentation.core_properties.title = "CareerCompass Graduation Project Presentation"
    presentation.core_properties.subject = (
        "AI-assisted graduate skill-gap analysis, personalized learning, and mentor guidance"
    )
    presentation.core_properties.author = "Basil Mohammad Alrmuny; Mohammed Al-Madhoun"
    presentation.core_properties.keywords = (
        "CareerCompass, graduation project, skill gap, career readiness, mentor guidance"
    )
    presentation.core_properties.comments = (
        "Formal defense deck generated from the current CareerCompass report and implementation."
    )

    with tempfile.TemporaryDirectory(prefix="careercompass-pptx-") as temporary:
        assets = extract_assets(REPORT, Path(temporary))

        # Slide 1 — cover
        slide = presentation.slides.add_slide(blank_layout)
        set_slide_background(slide, WHITE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.24, SLIDE_H, fill=NAVY, line=NAVY, line_width=0.1)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.24, 0, 0.08, SLIDE_H, fill=TEAL, line=TEAL, line_width=0.1)
        add_pill(slide, "GRADUATION PROJECT  •  2025/2026", 0.83, 0.58, 3.05, 0.40, size=10)
        add_text(
            slide,
            "CareerCompass",
            0.82,
            1.42,
            7.5,
            0.72,
            size=40,
            font_color=NAVY,
            bold=True,
            font_name=FONT_DISPLAY,
        )
        add_text(
            slide,
            "AI-Assisted Graduate Skill-Gap Analysis,\nPersonalized Learning, and Mentor Guidance",
            0.84,
            2.20,
            7.9,
            1.10,
            size=23,
            font_color=TEAL,
            bold=False,
            line_spacing=1.03,
        )
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.84, 3.48, 1.05, 0.06, fill=TEAL, line=TEAL, line_width=0.1)
        add_rich_line(
            slide,
            [
                ("Basil Mohammad Alrmuny", 14, NAVY, True),
                ("  •  202410100", 13, MUTED, False),
            ],
            0.84,
            3.83,
            5.2,
            0.34,
        )
        add_rich_line(
            slide,
            [
                ("Mohammed Al-Madhoun", 14, NAVY, True),
                ("  •  202411766", 13, MUTED, False),
            ],
            0.84,
            4.24,
            5.2,
            0.34,
        )
        add_rich_line(
            slide,
            [
                ("Supervisor: ", 12.5, MUTED, False),
                ("Dr. Shadi Ettantawi", 12.5, NAVY, True),
            ],
            0.84,
            4.86,
            5.2,
            0.34,
        )
        add_image_contain(slide, assets["meu"], 9.52, 0.44, 2.95, 1.35)
        add_compass(slide, 9.48, 2.33, 2.75)
        add_text(
            slide,
            "Middle East University\nFaculty of Information Technology",
            8.75,
            5.65,
            4.0,
            0.74,
            size=12,
            font_color=MUTED,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "From academic evidence to actionable career development",
            0.84,
            6.55,
            7.2,
            0.38,
            size=12.5,
            font_color=MUTED,
        )

        # Slide 2 — problem
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 2, "Why CareerCompass?", "Problem and motivation")
        add_text(
            slide,
            "Graduates know what they studied—not necessarily what their target career expects.",
            0.66,
            1.36,
            8.05,
            0.55,
            size=21,
            font_color=NAVY,
            bold=True,
        )
        problem_cards = [
            (
                "Academic records are indirect",
                "Course names and grades do not immediately reveal market-relevant skills or evidence strength.",
                TEAL,
            ),
            (
                "Course choice is often generic",
                "Learners must diagnose their own gaps and search large catalogues without a prioritized path.",
                BLUE,
            ),
            (
                "Guidance is fragmented",
                "Assessment, learning resources, mentor advice, and employer workflows usually live in separate tools.",
                AMBER,
            ),
        ]
        for index, (title, body, accent) in enumerate(problem_cards):
            add_two_line_card(slide, title, body, 0.68, 2.04 + index * 1.37, 8.0, 1.12, accent=accent)
            add_number_badge(slide, str(index + 1), 0.88, 2.31 + index * 1.37, 0.48, fill=accent)
        add_card(slide, 9.12, 1.53, 3.53, 4.93, fill=NAVY, line=NAVY, shadow=False)
        add_text(slide, "CONTEXT", 9.48, 1.94, 2.75, 0.28, size=10.5, font_color="9ED9D2", bold=True)
        add_text(slide, "21.2%", 9.46, 2.48, 2.75, 0.75, size=38, font_color=WHITE, bold=True, font_name=FONT_DISPLAY)
        add_text(
            slide,
            "Jordanian unemployment\nin Q4 2025",
            9.48,
            3.24,
            2.75,
            0.72,
            size=16,
            font_color=WHITE,
            bold=True,
        )
        add_shape(slide, MSO_SHAPE.RECTANGLE, 9.48, 4.18, 0.66, 0.05, fill=TEAL_2, line=TEAL_2, line_width=0.1)
        add_text(
            slide,
            "This motivates the problem. It is not evidence that the prototype reduced unemployment.",
            9.48,
            4.48,
            2.70,
            1.05,
            size=12.5,
            font_color="DCE7EF",
        )
        add_text(
            slide,
            "Sources: CareerCompass Report, Chapter 1; Trading Economics (2026); Jordan Department of Statistics (2025).",
            0.68,
            6.67,
            11.8,
            0.24,
            size=8.5,
            font_color=MUTED,
        )

        # Slide 3 — knowledge gap
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 3, "The knowledge gap is integration", "Related work")
        add_text(
            slide,
            "Existing systems solve parts of the problem. CareerCompass connects them through one evidence-based skill profile.",
            0.66,
            1.37,
            11.9,
            0.46,
            size=18,
            font_color=MUTED,
        )
        add_card(slide, 0.68, 2.03, 4.23, 4.40, fill=WHITE)
        add_pill(slide, "FRAGMENTED TODAY", 0.98, 2.28, 1.65, 0.34, fill=SLATE_PALE, font_color=NAVY, size=9.5)
        fragmented = [
            ("Academic transcript", "Evidence with no market interpretation"),
            ("Course platforms", "Catalogue without personal diagnosis"),
            ("Job platforms", "Keywords and self-declared skills"),
            ("Mentor guidance", "Useful but disconnected from evidence"),
        ]
        for i, (title, body) in enumerate(fragmented):
            y = 2.82 + i * 0.83
            add_shape(slide, MSO_SHAPE.OVAL, 1.00, y + 0.05, 0.38, 0.38, fill=SLATE_PALE, line=BORDER, line_width=0.6)
            add_text(slide, str(i + 1), 1.00, y + 0.05, 0.38, 0.38, size=10, font_color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, title, 1.56, y, 2.78, 0.26, size=14, font_color=NAVY, bold=True)
            add_text(slide, body, 1.56, y + 0.29, 2.93, 0.34, size=10.5, font_color=MUTED)
        add_right_arrow(slide, 5.27, 3.56, 1.05, 0.58, fill=TEAL)
        add_card(slide, 6.70, 2.03, 5.93, 4.40, fill=TEAL_PALE, line="B9DBD6")
        add_pill(slide, "CAREERCOMPASS", 7.02, 2.28, 1.48, 0.34, fill=TEAL, font_color=WHITE, size=9.5)
        add_text(slide, "One shared Student Skill Vector", 7.02, 2.83, 5.15, 0.42, size=22, font_color=NAVY, bold=True)
        add_text(
            slide,
            "Confirmed academic evidence becomes the consistent input to every downstream decision.",
            7.02,
            3.29,
            4.95,
            0.55,
            size=13.5,
            font_color=MUTED,
        )
        integrated = [
            ("Skill-gap dashboard", TEAL),
            ("Personalized learning", BLUE),
            ("Quiz-based feedback", AMBER),
            ("Mentor guidance", GREEN),
        ]
        for i, (label, accent) in enumerate(integrated):
            x = 7.02 + (i % 2) * 2.55
            y = 4.08 + (i // 2) * 0.88
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 2.28, 0.62, fill=WHITE, line=accent, line_width=1.0)
            add_shape(slide, MSO_SHAPE.OVAL, x + 0.15, y + 0.15, 0.31, 0.31, fill=accent, line=accent, line_width=0.2)
            add_text(slide, label, x + 0.56, y + 0.10, 1.55, 0.42, size=11.5, font_color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(
            slide,
            "Research contribution: an integrated academic prototype—not a claim of measured employment impact.",
            0.68,
            6.67,
            11.9,
            0.26,
            size=9.5,
            font_color=MUTED,
        )

        # Slide 4 — objective and scope
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 4, "Objective and project scope", "What the system is designed to achieve")
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.68, 1.40, 11.98, 0.72, fill=NAVY, line=NAVY, line_width=0.2)
        add_text(
            slide,
            "Convert confirmed academic evidence into an explainable plan for career readiness.",
            1.02,
            1.50,
            11.25,
            0.45,
            size=20,
            font_color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        pillars = [
            ("Analyze", "Parse and confirm courses and grades.", TEAL),
            ("Identify", "Measure proficiency against a chosen career path.", BLUE),
            ("Improve", "Prioritize gaps, recommend courses, and assess progress.", AMBER),
            ("Guide", "Connect students with mentors for human support.", GREEN),
        ]
        for i, (title, body, accent) in enumerate(pillars):
            x = 0.68 + i * 3.03
            add_card(slide, x, 2.50, 2.79, 2.48, fill=WHITE)
            add_number_badge(slide, str(i + 1), x + 0.20, 2.72, 0.50, fill=accent)
            add_text(slide, title, x + 0.20, 3.37, 2.38, 0.34, size=18, font_color=NAVY, bold=True)
            add_text(slide, body, x + 0.20, 3.81, 2.35, 0.80, size=12.2, font_color=MUTED)
        add_card(slide, 0.68, 5.35, 11.98, 1.20, fill=AMBER_PALE, line="E8C98A", shadow=False)
        add_pill(slide, "SCOPE BOUNDARY", 0.96, 5.59, 1.48, 0.32, fill=AMBER, font_color=WHITE, size=9.2)
        add_rich_line(
            slide,
            [
                ("Current release: ", 13.5, NAVY, True),
                ("student analysis, learning support, and mentor workflows.   ", 13.5, INK, False),
                ("Deferred [L]: ", 13.5, RED, True),
                ("real AI job/candidate ranking.", 13.5, INK, False),
            ],
            2.72,
            5.57,
            9.25,
            0.42,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            "Job posting and vacancy browsing are implemented; development-mock scores are not presented as AI results.",
            2.72,
            6.05,
            9.20,
            0.27,
            size=10.5,
            font_color=MUTED,
        )

        # Slide 5 — actors and journey
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 5, "Five actors, one governed workflow", "Users and end-to-end journey")
        actors = [
            ("JS", "Job seeker", "evidence & learning"),
            ("CM", "Content manager", "reviewed course maps"),
            ("EX", "Expert", "mentor guidance"),
            ("EMP", "Employer", "vacancy management"),
            ("ADM", "Administrator", "governance & access"),
        ]
        for i, (abbr, title, detail) in enumerate(actors):
            x = 0.65 + i * 2.54
            add_card(slide, x, 1.42, 2.30, 0.95, fill=WHITE, shadow=False)
            add_shape(slide, MSO_SHAPE.OVAL, x + 0.16, 1.61, 0.48, 0.48, fill=TEAL_PALE, line=TEAL, line_width=0.7)
            add_text(slide, abbr, x + 0.16, 1.61, 0.48, 0.48, size=9.5, font_color=TEAL, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, title, x + 0.75, 1.53, 1.37, 0.24, size=11.5, font_color=NAVY, bold=True)
            add_text(slide, detail, x + 0.75, 1.81, 1.37, 0.26, size=9.2, font_color=MUTED)
        add_text(slide, "CORE STUDENT JOURNEY", 0.68, 2.86, 3.2, 0.26, size=10, font_color=TEAL, bold=True)
        journey = [
            ("Select\ncareer path", "1"),
            ("Upload &\nconfirm", "2"),
            ("Build skill\nvector", "3"),
            ("Prioritize\ngaps", "4"),
            ("Learn &\nassess", "5"),
            ("Book a\nmentor", "6"),
        ]
        for i, (label, num) in enumerate(journey):
            x = 0.68 + i * 2.04
            shape = add_shape(slide, MSO_SHAPE.CHEVRON, x, 3.30, 2.18, 1.02, fill=WHITE if i % 2 == 0 else TEAL_PALE, line="B9D5D1", line_width=0.8)
            add_shape(slide, MSO_SHAPE.OVAL, x + 0.12, 3.39, 0.40, 0.40, fill=TEAL, line=TEAL, line_width=0.2)
            add_text(slide, num, x + 0.12, 3.39, 0.40, 0.40, size=9.5, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, label, x + 0.59, 3.42, 1.16, 0.60, size=11.5, font_color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_card(slide, 0.68, 4.82, 11.98, 1.38, fill=WHITE, shadow=False)
        add_text(slide, "Governance around the journey", 0.95, 5.08, 3.45, 0.30, size=16, font_color=NAVY, bold=True)
        governance = [
            "Content managers approve extracted syllabus skills before publication.",
            "Experts see authorized mentee data only through accepted consultations.",
            "Employers manage vacancies; real match scoring remains outside the current release.",
        ]
        add_bullets(slide, governance, 4.14, 5.03, 7.95, 0.95, size=11.5, gap=4)

        # Slide 6 — architecture
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 6, "System architecture", "Docker Compose deployment view")
        add_card(slide, 0.68, 1.37, 11.98, 5.58, fill=WHITE, shadow=False)
        add_image_contain(slide, assets["architecture"], 0.82, 1.47, 11.70, 5.38)

        # Slide 7 — data and pipeline
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 7, "From data to decisions", "Knowledge base and analysis pipeline")
        metrics = [
            ("903", "canonical skills", "versioned taxonomy"),
            ("771", "career requirements", "across 9 career paths"),
            ("2,238", "job postings", "market-demand evidence"),
            ("20 / 114", "real syllabi", "current coverage limit"),
        ]
        for i, (value, label, note) in enumerate(metrics):
            add_metric_card(slide, value, label, 0.67 + i * 3.04, 1.38, 2.78, note=note)
        groups = [
            ("1  INPUT EVIDENCE", ["Confirmed transcript", "Reviewed course outcomes", "Career-path ontology"], 0.68, BLUE_PALE, BLUE),
            ("2  DETERMINISTIC CORE", ["Normalize skills", "Build skill vector", "Compute gap × demand"], 4.55, TEAL_PALE, TEAL),
            ("3  GUIDANCE OUTPUTS", ["Dashboard & explanations", "Courses and quizzes", "Mentor ranking"], 8.42, AMBER_PALE, AMBER),
        ]
        for title, items, x, fill, accent in groups:
            add_card(slide, x, 3.05, 3.42, 2.80, fill=fill, line=accent, shadow=False)
            add_text(slide, title, x + 0.24, 3.30, 2.94, 0.30, size=11.5, font_color=accent, bold=True)
            for i, item in enumerate(items):
                add_shape(slide, MSO_SHAPE.OVAL, x + 0.26, 3.90 + i * 0.57, 0.29, 0.29, fill=accent, line=accent, line_width=0.2)
                add_text(slide, item, x + 0.70, 3.82 + i * 0.57, 2.33, 0.43, size=12.4, font_color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_right_arrow(slide, 4.08, 4.11, 0.38, 0.52, fill=TEAL)
        add_right_arrow(slide, 7.95, 4.11, 0.38, 0.52, fill=AMBER)
        add_card(slide, 1.78, 6.18, 9.78, 0.60, fill=WHITE, shadow=False)
        add_pill(slide, "LLM BOUNDARY", 2.04, 6.31, 1.30, 0.30, fill=NAVY, font_color=WHITE, size=8.8)
        add_text(
            slide,
            "Generation is selective: an LLM may create quiz questions or prose, but it cannot change stored proficiency, gaps, priorities, or grading arithmetic.",
            3.62,
            6.24,
            7.47,
            0.40,
            size=10.5,
            font_color=MUTED,
            valign=MSO_ANCHOR.MIDDLE,
        )

        # Slide 8 — deterministic core
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 8, "Explainable skill-gap prioritization", "Deterministic analytical core")
        add_pill(slide, "NO MODEL CALL IN SCORING", 0.68, 1.40, 2.18, 0.36, fill=GREEN_PALE, font_color=GREEN, size=9.5)
        formula_cards = [
            ("Student proficiency", "grade-weighted attainment\n÷ graded evidence", BLUE),
            ("Skill gap", "max(target − current, 0)", TEAL),
            ("Action priority", "skill gap × market demand", AMBER),
        ]
        for i, (title, formula, accent) in enumerate(formula_cards):
            y = 2.02 + i * 1.30
            add_card(slide, 0.68, y, 3.92, 1.05, fill=WHITE, shadow=False)
            add_shape(slide, MSO_SHAPE.RECTANGLE, 0.68, y, 0.07, 1.05, fill=accent, line=accent, line_width=0.1)
            add_text(slide, title, 0.96, y + 0.17, 3.26, 0.28, size=14, font_color=NAVY, bold=True)
            add_text(slide, formula, 0.96, y + 0.53, 3.26, 0.37, size=12.2, font_color=accent, bold=True)
        add_card(slide, 4.88, 1.42, 7.78, 4.83, fill=WHITE, shadow=False)
        add_image_contain(slide, assets["gap_algorithm"], 5.06, 1.63, 7.42, 4.46)
        add_card(slide, 0.68, 6.14, 11.98, 0.58, fill=TEAL_PALE, line="B9DBD6", shadow=False)
        add_text(
            slide,
            "The LLM may explain the computed result in plain language; validation prevents it from altering any numeric classification or priority.",
            0.96,
            6.24,
            11.42,
            0.33,
            size=11.2,
            font_color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )

        # Slide 9 — UI
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 9, "The implemented student experience", "Prototype walkthrough")
        ui_cards = [
            ("1", "Confirm academic evidence", "Review extracted courses and grades before anything is saved.", assets["transcript"]),
            ("2", "See readiness and gaps", "Understand evidence coverage, strengths, and prioritized missing skills.", assets["dashboard"]),
            ("3", "Request mentor guidance", "Browse ranked mentors and book within published availability.", assets["mentors"]),
        ]
        for i, (num, title, body, image_path) in enumerate(ui_cards):
            x = 0.57 + i * 4.25
            add_card(slide, x, 1.45, 3.95, 4.92, fill=WHITE)
            add_image_cover(slide, image_path, x + 0.15, 1.63, 3.65, 2.30)
            add_picture_frame(slide, x + 0.15, 1.63, 3.65, 2.30)
            add_number_badge(slide, num, x + 0.22, 4.20, 0.48, fill=TEAL)
            add_text(slide, title, x + 0.82, 4.18, 2.82, 0.36, size=15, font_color=NAVY, bold=True)
            add_text(slide, body, x + 0.22, 4.77, 3.47, 0.90, size=11.4, font_color=MUTED)
            add_pill(slide, "CURRENT UI", x + 0.22, 5.83, 1.02, 0.30, fill=TEAL_PALE, font_color=TEAL, size=8.5)
        add_text(
            slide,
            "Screenshots show the current demonstration build. Synthetic-data warnings remain visible by design.",
            0.68,
            6.68,
            11.9,
            0.24,
            size=9.5,
            font_color=MUTED,
        )

        # Slide 10 — stack and security
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 10, "Engineering design and security controls", "Implementation")
        columns = [
            (
                "Frontend",
                "React 19 + TypeScript",
                ["Role-protected navigation", "React Query data flows", "Responsive, accessible states"],
                BLUE,
                BLUE_PALE,
            ),
            (
                "Business API",
                "Java 17 + Spring Boot",
                ["Authentication and RBAC", "Transactional workflows", "JPA, Flyway, OpenAPI"],
                TEAL,
                TEAL_PALE,
            ),
            (
                "Analysis service",
                "Python + FastAPI",
                ["Parsing and skill analysis", "Deterministic ranking core", "PostgreSQL knowledge data"],
                AMBER,
                AMBER_PALE,
            ),
        ]
        for i, (title, stack, bullets, accent, fill) in enumerate(columns):
            x = 0.68 + i * 4.05
            add_card(slide, x, 1.50, 3.68, 3.66, fill=WHITE)
            add_shape(slide, MSO_SHAPE.RECTANGLE, x, 1.50, 3.68, 0.12, fill=accent, line=accent, line_width=0.1)
            add_pill(slide, title.upper(), x + 0.25, 1.90, 1.34, 0.32, fill=fill, font_color=accent, size=8.8)
            add_text(slide, stack, x + 0.25, 2.44, 3.15, 0.42, size=19, font_color=NAVY, bold=True)
            add_bullets(slide, bullets, x + 0.25, 3.08, 3.08, 1.62, size=12.4, bullet_color=accent, gap=8)
        add_card(slide, 0.68, 5.52, 11.98, 1.10, fill=NAVY, line=NAVY, shadow=False)
        add_text(slide, "DEFENSE IN DEPTH", 0.98, 5.73, 1.72, 0.28, size=10, font_color="9ED9D2", bold=True)
        controls = ["BCrypt passwords", "JWT + revocation", "Role ownership checks", "Service bearer token", "Bounded PDF validation"]
        for i, item in enumerate(controls):
            x = 2.85 + i * 1.83
            add_shape(slide, MSO_SHAPE.OVAL, x, 5.74, 0.29, 0.29, fill=TEAL_2, line=TEAL_2, line_width=0.2)
            add_text(slide, item, x + 0.38, 5.66, 1.34, 0.50, size=10.4, font_color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)

        # Slide 11 — testing
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 11, "Verification across layers", "Reported test evidence")
        layers = [
            ("Database constraints", "103 operations", NAVY, 5.0),
            ("Backend & actor workflows", "219 test methods reported", TEAL, 4.55),
            ("AI algorithm modules", "808 assertion checks / 216 tests", BLUE, 4.10),
            ("HTTP boundary & adversarial", "19 routes / 61 malformed cases", AMBER, 3.65),
        ]
        for i, (title, metric, accent, width) in enumerate(layers):
            x = 0.78 + (5.15 - width) / 2
            y = 1.55 + i * 1.04
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, 0.82, fill=WHITE, line=accent, line_width=1.0)
            add_text(slide, title, x + 0.24, y + 0.13, width - 0.48, 0.25, size=13.2, font_color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, metric, x + 0.24, y + 0.45, width - 0.48, 0.20, size=10.5, font_color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_card(slide, 6.25, 1.55, 6.36, 4.63, fill=WHITE)
        add_text(slide, "What the evidence demonstrates", 6.58, 1.88, 5.75, 0.36, size=18, font_color=NAVY, bold=True)
        evidence = [
            "Six live Java–FastAPI contract tests are documented.",
            "Repeated vector and gap calls produced byte-identical JSON.",
            "58 of 61 adversarial inputs initially returned the expected 4xx response; three defects were fixed.",
            "Critical and high-severity defects in the recorded audits were remediated.",
        ]
        add_bullets(slide, evidence, 6.60, 2.47, 5.43, 2.75, size=13.2, gap=9)
        add_card(slide, 6.58, 5.40, 5.65, 0.55, fill=AMBER_PALE, line="E8C98A", shadow=False)
        add_text(slide, "Evidence supports the prototype; it does not prove every NFR or real-world impact.", 6.82, 5.50, 5.16, 0.29, size=10.2, font_color=AMBER, bold=True, align=PP_ALIGN.CENTER)
        add_text(
            slide,
            "Source: current report, Chapter 7 and Appendices C–G. Counts are presented as documented evidence, not independently re-executed acceptance results.",
            0.68,
            6.66,
            11.9,
            0.28,
            size=8.8,
            font_color=MUTED,
        )

        # Slide 12 — achieved vs deferred
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 12, "Objective status: substantially achieved", "Honest scope assessment")
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.68, 1.40, 11.98, 0.67, fill=TEAL_PALE, line="B9DBD6", line_width=0.8)
        add_text(slide, "The career-readiness workflow works end to end; the original job-matching objective is not complete.", 0.98, 1.52, 11.37, 0.38, size=18, font_color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_card(slide, 0.68, 2.36, 5.74, 3.94, fill=GREEN_PALE, line="B8DCC9", shadow=False)
        add_pill(slide, "ACHIEVED", 0.98, 2.67, 1.13, 0.34, fill=GREEN, font_color=WHITE, size=9.4)
        achieved = [
            "Transcript upload, review, and confirmation",
            "Skill vector and prioritized gap dashboard",
            "Catalogue-based recommendations; quizzes when an LLM is configured",
            "Deterministic mentor ranking and consultation booking",
            "Five-role authentication, authorization, and administration",
        ]
        add_bullets(slide, achieved, 0.98, 3.24, 5.03, 2.62, size=13.2, bullet_color=GREEN, gap=8)
        add_card(slide, 6.91, 2.36, 5.75, 3.94, fill=AMBER_PALE, line="E8C98A", shadow=False)
        add_pill(slide, "DEFERRED OR LIMITED", 7.21, 2.67, 1.95, 0.34, fill=AMBER, font_color=WHITE, size=9.2)
        limited = [
            "Real AI job/candidate scoring returns HTTP 501 — Low priority",
            "Only 20 of 114 course syllabi are available as real evidence",
            "Course-catalogue coverage remains partial",
            "Accuracy, fairness, usability, and employment impact lack large-scale evaluation",
            "Single-host demonstration is not a production deployment",
        ]
        add_bullets(slide, limited, 7.21, 3.24, 4.98, 2.65, size=12.9, bullet_color=AMBER, gap=7)

        # Slide 13 — roadmap
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 13, "Limitations translated into a roadmap", "Future work by priority")
        roadmap = [
            (
                "HIGH",
                "Evidence quality",
                ["Collect and approve remaining real syllabi", "Run human-labelled quality and relevance evaluation", "Close security, accessibility, and CI evidence gaps"],
                GREEN,
                GREEN_PALE,
            ),
            (
                "MEDIUM",
                "Operational maturity",
                ["Cache catalogues and indexes; add observability", "Expand Arabic and CV support", "Collect explicit mentor expertise signals"],
                BLUE,
                BLUE_PALE,
            ),
            (
                "LOW",
                "Deferred matching",
                ["Implement real candidate scoring", "Evaluate against human-reviewed rankings", "Keep mock scoring restricted to tests"],
                AMBER,
                AMBER_PALE,
            ),
        ]
        for i, (priority, title, bullets, accent, fill) in enumerate(roadmap):
            x = 0.68 + i * 4.04
            add_card(slide, x, 1.55, 3.66, 4.72, fill=WHITE)
            add_pill(slide, f"{priority} PRIORITY", x + 0.26, 1.85, 1.46, 0.34, fill=accent, font_color=WHITE, size=9.2)
            add_text(slide, title, x + 0.26, 2.45, 3.04, 0.40, size=19, font_color=NAVY, bold=True)
            add_bullets(slide, bullets, x + 0.26, 3.10, 3.04, 2.15, size=12.7, bullet_color=accent, gap=10)
            add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.26, 5.65, 3.03, 0.08, fill=fill, line=fill, line_width=0.1)
        add_card(slide, 3.02, 6.48, 7.30, 0.40, fill=AMBER_PALE, line=AMBER_PALE, shadow=False)
        add_text(slide, "Real AI job matching remains explicitly [L] in the functional requirements and future-work plan.", 3.19, 6.51, 6.96, 0.27, size=10.5, font_color=AMBER, bold=True, align=PP_ALIGN.CENTER)

        # Slide 14 — closing
        slide = presentation.slides.add_slide(blank_layout)
        set_slide_background(slide, WHITE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.22, SLIDE_H, fill=NAVY, line=NAVY, line_width=0.1)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.22, 0, 0.07, SLIDE_H, fill=TEAL, line=TEAL, line_width=0.1)
        add_pill(slide, "CONCLUSION", 0.82, 0.58, 1.28, 0.36, fill=TEAL_PALE, font_color=TEAL, size=9.4)
        add_text(
            slide,
            "CareerCompass turns academic evidence\ninto actionable career development.",
            0.82,
            1.36,
            7.64,
            1.16,
            size=31,
            font_color=NAVY,
            bold=True,
            font_name=FONT_DISPLAY,
        )
        takeaways = [
            ("Evidence-based", "Starts with confirmed courses and grades."),
            ("Explainable", "Keeps scoring deterministic and traceable."),
            ("Human-centered", "Combines learning support with mentor guidance."),
        ]
        for i, (title, body) in enumerate(takeaways):
            x = 0.82 + i * 2.66
            add_card(slide, x, 3.12, 2.42, 1.40, fill=CANVAS, shadow=False)
            add_shape(slide, MSO_SHAPE.OVAL, x + 0.18, 3.38, 0.34, 0.34, fill=TEAL, line=TEAL, line_width=0.2)
            add_text(slide, title, x + 0.65, 3.29, 1.52, 0.30, size=13.5, font_color=NAVY, bold=True)
            add_text(slide, body, x + 0.18, 3.78, 2.03, 0.50, size=10.4, font_color=MUTED)
        add_text(slide, "Questions & Discussion", 0.82, 5.23, 5.80, 0.58, size=27, font_color=TEAL, bold=True, font_name=FONT_DISPLAY)
        add_text(slide, "Basil Mohammad Alrmuny  •  Mohammed Al-Madhoun", 0.84, 5.95, 6.80, 0.34, size=12.5, font_color=MUTED)
        add_compass(slide, 9.63, 2.03, 2.53, muted=True)
        add_image_contain(slide, assets["meu"], 9.48, 5.42, 3.02, 1.25)
        add_text(slide, "Main contribution: a validated academic prototype for skill-gap analysis, personalized learning, and mentor guidance.", 0.84, 6.57, 7.75, 0.40, size=10.5, font_color=MUTED)
        add_footer(slide, 14)

        # Appendix 1 — skill vector
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 15, "Student Skill Vector construction", "Appendix A1", appendix=True)
        add_card(slide, 0.68, 1.43, 11.98, 4.93, fill=WHITE, shadow=False)
        add_image_contain(slide, assets["skill_vector"], 0.88, 1.63, 11.58, 4.52)
        add_text(slide, "Grades estimate proficiency; evidence coverage remains separate; a quiz score can replace the grade-derived estimate without erasing its origin.", 0.82, 6.52, 11.66, 0.34, size=10.5, font_color=MUTED, align=PP_ALIGN.CENTER)

        # Appendix 2 — course and quiz
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 16, "Learning recommendation and assessment", "Appendix A2", appendix=True)
        add_card(slide, 0.62, 1.42, 6.00, 4.95, fill=WHITE, shadow=False)
        add_card(slide, 6.72, 1.42, 6.00, 4.95, fill=WHITE, shadow=False)
        add_image_contain(slide, assets["course_ranking"], 0.82, 1.64, 5.60, 3.92)
        add_image_contain(slide, assets["quiz"], 6.92, 1.64, 5.60, 3.92)
        add_text(slide, "Deterministic course ranking", 1.12, 5.72, 5.00, 0.32, size=15, font_color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "LLM generation; programmatic validation and grading", 7.03, 5.72, 5.42, 0.32, size=14.4, font_color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "Recommendation output never invents a catalogue item; invalid quiz questions are rejected rather than used.", 1.72, 6.53, 9.89, 0.32, size=10.5, font_color=MUTED, align=PP_ALIGN.CENTER)

        # Appendix 3 — requirements priorities
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 17, "Functional requirement priorities", "Appendix A3", appendix=True)
        add_text(slide, "112 functional requirements are explicitly prioritized for release planning.", 0.68, 1.40, 11.95, 0.40, size=18, font_color=MUTED)
        total = 112
        segments = [(77, "HIGH", GREEN), (30, "MEDIUM", BLUE), (5, "LOW", AMBER)]
        bar_x, bar_y, bar_w, bar_h = 0.83, 2.24, 11.67, 0.78
        cursor = bar_x
        for count, label, accent in segments:
            width = bar_w * count / total
            add_shape(slide, MSO_SHAPE.RECTANGLE, cursor, bar_y, width, bar_h, fill=accent, line=accent, line_width=0.2)
            add_text(slide, f"{count} {label}", cursor, bar_y, width, bar_h, size=13 if width > 1.5 else 10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            cursor += width
        add_text(slide, "68.8%", 0.83, 3.17, 2.15, 0.42, size=22, font_color=GREEN, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "High priority", 0.83, 3.61, 2.15, 0.28, size=11, font_color=MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, "26.8%", 5.58, 3.17, 2.15, 0.42, size=22, font_color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "Medium priority", 5.58, 3.61, 2.15, 0.28, size=11, font_color=MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, "4.5%", 10.35, 3.17, 2.15, 0.42, size=22, font_color=AMBER, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "Low priority", 10.35, 3.61, 2.15, 0.28, size=11, font_color=MUTED, align=PP_ALIGN.CENTER)
        add_card(slide, 1.40, 4.30, 10.53, 1.62, fill=AMBER_PALE, line="E8C98A", shadow=False)
        add_pill(slide, "REAL AI JOB MATCHING — [L]", 1.76, 4.60, 2.35, 0.34, fill=AMBER, font_color=WHITE, size=9.0)
        add_text(slide, "FR-JS-23   •   FR-EMP-11   •   FR-EMP-19   •   FR-AI-12   •   FR-AI-13", 1.77, 5.13, 9.78, 0.34, size=15, font_color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "Priority indicates implementation/release importance—not completion status.", 2.90, 6.30, 7.56, 0.30, size=10.5, font_color=MUTED, align=PP_ALIGN.CENTER)

        # Appendix 4 — references
        slide = presentation.slides.add_slide(blank_layout)
        add_slide_title(slide, 18, "Selected references and evidence sources", "Appendix A4", appendix=True)
        references_left = [
            "CareerCompass Graduation Project Report, Chapters 1, 3, 6–8 and Appendices C–G.",
            "CareerCompass repository: current React, Spring Boot, FastAPI, contracts, and implementation-status documents.",
            "Walker, R. E. (2024). Mapping curricula to skills and occupations using course descriptions. IEEE EDUNINE.",
            "Reimers, N., & Gurevych, I. (2019). Sentence-BERT. EMNLP-IJCNLP.",
        ]
        references_right = [
            "Beutling, A., & Spahic, A. (2024). Knowledge-based course recommendation aligned with career goals.",
            "Nouib, H. et al. (2025). Predicting graduate employability using hybrid AHP–TOPSIS and machine learning. Technologies, 13(9).",
            "Jordan Department of Statistics (2025). Unemployment indicators.",
            "Trading Economics (2026). Jordan unemployment rate; historical 2025 figures cited in the report.",
        ]
        add_card(slide, 0.68, 1.48, 5.78, 4.95, fill=WHITE, shadow=False)
        add_card(slide, 6.86, 1.48, 5.78, 4.95, fill=WHITE, shadow=False)
        add_bullets(slide, references_left, 1.02, 1.88, 5.08, 4.18, size=12.1, gap=13)
        add_bullets(slide, references_right, 7.20, 1.88, 5.08, 4.18, size=12.1, gap=13)
        add_text(slide, "External statistics motivate the problem; implementation and test claims are grounded in the submitted report and repository.", 1.22, 6.60, 10.90, 0.30, size=10.2, font_color=MUTED, align=PP_ALIGN.CENTER)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "output",
        nargs="?",
        type=Path,
        default=ROOT / "CareerCompass_Graduation_Project_Presentation.pptx",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    make_deck(args.output)
    print(f"Created {args.output}")


if __name__ == "__main__":
    main()
